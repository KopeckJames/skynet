# processors/file_processor.py
from pathlib import Path
from typing import Dict, Any
import tempfile
import fitz  # PyMuPDF
import pdf2image
import pytesseract
import pandas as pd
from pptx import Presentation
from .base_processor import BaseProcessor

class FileProcessor(BaseProcessor):
    """Processor for handling various document file types"""

    def process(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """Process file based on its type"""
        if not self.validate_content(content):
            raise ValueError("Invalid content provided")

        file_path = content.get("path")
        file_type = self.get_file_type(file_path)
        
        try:
            if file_type == "application/pdf":
                processed_content = self.process_pdf(file_path)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                processed_content = self.process_powerpoint(file_path)
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                processed_content = self.process_excel(file_path)
            elif file_type == "text/csv":
                processed_content = self.process_csv(file_path)
            else:
                processed_content = self.process_text_file(file_path)

            return {
                "content": processed_content,
                "file_type": file_type,
                "metadata": {
                    "filename": Path(file_path).name,
                    "file_type": file_type,
                    "file_size": Path(file_path).stat().st_size
                }
            }
        except Exception as e:
            raise Exception(f"Error processing file: {str(e)}")

    def process_pdf(self, file_path: str) -> str:
        """Extract text from PDF files with OCR fallback"""
        try:
            text = []
            doc = fitz.open(file_path)
            
            # First try normal text extraction
            for page in doc:
                text.append(page.get_text())
            
            content = "\n".join(text)
            
            # If no text found, try OCR
            if not content.strip():
                images = pdf2image.convert_from_path(file_path)
                ocr_text = []
                for image in images:
                    ocr_text.append(pytesseract.image_to_string(image))
                content = "\n".join(ocr_text)
            
            return content.strip()
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")

    def process_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")

    def process_excel(self, file_path: str) -> str:
        """Process Excel files"""
        try:
            # Read all sheets
            all_sheets = pd.read_excel(file_path, sheet_name=None)
            text_content = []
            
            for sheet_name, df in all_sheets.items():
                text_content.append(f"Sheet: {sheet_name}")
                text_content.append(df.to_string())
                text_content.append("\n")
            
            return "\n".join(text_content)
        except Exception as e:
            raise Exception(f"Error processing Excel: {str(e)}")

    def process_csv(self, file_path: str) -> str:
        """Process CSV files"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            raise Exception(f"Error processing CSV: {str(e)}")

    def process_text_file(self, file_path: str) -> str:
        """Process plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read().strip()
        except UnicodeDecodeError:
            # Try different encodings if UTF-8 fails
            encodings = ['latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read().strip()
                except UnicodeDecodeError:
                    continue
            raise Exception("Unable to decode file with any supported encoding")

    @staticmethod
    def get_file_type(filename: str) -> str:
        """Determine file type based on extension"""
        ext = Path(filename).suffix.lower()
        
        # Map of extensions to MIME types
        mime_types = {
            '.txt': 'text/plain',
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls': 'application/vnd.ms-excel',
            '.csv': 'text/csv',
            '.doc': 'application/msword',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.md': 'text/markdown',
            '.json': 'application/json',
            '.xml': 'application/xml',
            '.html': 'text/html',
            '.htm': 'text/html'
        }
        
        return mime_types.get(ext, 'application/octet-stream')

    def clean_extracted_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
            
        # Remove excessive whitespace
        text = " ".join(text.split())
        
        # Remove common artifacts
        artifacts = ['\x00', '\x0c', '\x1a']
        for artifact in artifacts:
            text = text.replace(artifact, ' ')
            
        return text.strip()